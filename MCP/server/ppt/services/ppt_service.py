"""
PowerPoint Service Module

This module provides the core PowerPoint manipulation functionality,
separate from the MCP server layer to follow proper separation of concerns.
"""

from typing import List, Dict, Any, Optional
from pathlib import Path
import logging
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)


class PresentationManager:
    """Manages PowerPoint presentation creation and manipulation."""
    
    def __init__(self):
        self.presentations: Dict[str, Presentation] = {}
    
    def create_presentation(self, presentation_id: str, template_path: Optional[str] = None) -> Dict[str, Any]:
        """
        Create a new PowerPoint presentation.
        
        Args:
            presentation_id: Unique identifier for this presentation
            template_path: Optional path to a template .pptx file
            
        Returns:
            Dict with status and message
        """
        if presentation_id in self.presentations:
            return {
                "success": False,
                "message": f"Presentation '{presentation_id}' already exists"
            }
        
        if template_path and Path(template_path).exists():
            prs = Presentation(template_path)
            logger.info(f"Created presentation from template: {template_path}")
        else:
            prs = Presentation()
            logger.info(f"Created blank presentation")
        
        self.presentations[presentation_id] = prs
        
        return {
            "success": True,
            "message": f"Created presentation '{presentation_id}' with {len(prs.slides)} slides",
            "slide_count": len(prs.slides)
        }
    
    def add_title_slide(self, presentation_id: str, title: str, subtitle: str = "") -> Dict[str, Any]:
        """
        Add a title slide to the presentation.
        
        Args:
            presentation_id: Presentation identifier
            title: Main title text
            subtitle: Optional subtitle text
            
        Returns:
            Dict with status and message
        """
        if presentation_id not in self.presentations:
            return {"success": False, "message": f"Presentation '{presentation_id}' not found"}
        
        prs = self.presentations[presentation_id]
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = title
        if subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle
        
        return {
            "success": True,
            "message": f"Added title slide: {title}",
            "slide_number": len(prs.slides)
        }
    
    def add_content_slide(self, presentation_id: str, title: str, content: List[str],
                         font_size: Optional[int] = None, font_color: Optional[List[int]] = None,
                         title_font_size: Optional[int] = None, title_font_color: Optional[List[int]] = None,
                         border_color: Optional[List[int]] = None, border_width: Optional[float] = None) -> Dict[str, Any]:
        """
        Add a content slide with bullet points and optional styling.
        
        Args:
            presentation_id: Presentation identifier
            title: Slide title
            content: List of bullet points
            font_size: Optional font size in points (e.g., 18, 24)
            font_color: Optional font color as RGB list [R, G, B] (e.g., [0, 0, 0] for black)
            title_font_size: Optional title font size in points
            title_font_color: Optional title font color as RGB list
            border_color: Optional border color as RGB list
            border_width: Optional border width in points (e.g., 1.5)
            
        Returns:
            Dict with status and message
        """
        if presentation_id not in self.presentations:
            return {"success": False, "message": f"Presentation '{presentation_id}' not found"}
        
        prs = self.presentations[presentation_id]
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = title
        
        # Apply title styling if provided
        if title_font_size or title_font_color:
            title_frame = slide.shapes.title.text_frame
            for paragraph in title_frame.paragraphs:
                if title_font_size:
                    paragraph.font.size = Pt(title_font_size)
                if title_font_color and len(title_font_color) == 3:
                    paragraph.font.color.rgb = RGBColor(*title_font_color)
        
        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()
        
        for item in content:
            p = text_frame.add_paragraph()
            p.text = item
            p.level = 0
            
            # Apply content styling if provided
            if font_size:
                p.font.size = Pt(font_size)
            if font_color and len(font_color) == 3:
                p.font.color.rgb = RGBColor(*font_color)
        
        # Apply border to content text box if provided
        if border_color or border_width:
            shape = slide.placeholders[1]
            line = shape.line
            line.color.rgb = RGBColor(*border_color) if border_color and len(border_color) == 3 else RGBColor(0, 0, 0)
            line.width = Pt(border_width) if border_width else Pt(1.5)
        
        return {
            "success": True,
            "message": f"Added content slide: {title}",
            "slide_number": len(prs.slides),
            "content_items": len(content)
        }
    
    def add_two_column_slide(self, presentation_id: str, title: str, 
                           left_content: List[str], right_content: List[str],
                           font_size: Optional[int] = None, font_color: Optional[List[int]] = None,
                           title_font_size: Optional[int] = None, title_font_color: Optional[List[int]] = None,
                           border_color: Optional[List[int]] = None, border_width: Optional[float] = None) -> Dict[str, Any]:
        """
        Add a slide with two columns of content and optional styling.
        
        Args:
            presentation_id: Presentation identifier
            title: Slide title
            left_content: Left column bullet points
            right_content: Right column bullet points
            font_size: Optional font size in points (e.g., 18, 24)
            font_color: Optional font color as RGB list [R, G, B]
            title_font_size: Optional title font size in points
            title_font_color: Optional title font color as RGB list
            border_color: Optional border color as RGB list
            border_width: Optional border width in points
            
        Returns:
            Dict with status and message
        """
        if presentation_id not in self.presentations:
            return {"success": False, "message": f"Presentation '{presentation_id}' not found"}
        
        prs = self.presentations[presentation_id]
        slide_layout = prs.slide_layouts[3]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = title
        
        # Apply title styling if provided
        if title_font_size or title_font_color:
            title_frame = slide.shapes.title.text_frame
            for paragraph in title_frame.paragraphs:
                if title_font_size:
                    paragraph.font.size = Pt(title_font_size)
                if title_font_color and len(title_font_color) == 3:
                    paragraph.font.color.rgb = RGBColor(*title_font_color)
        
        # Left column
        left_frame = slide.placeholders[1].text_frame
        left_frame.clear()
        for item in left_content:
            p = left_frame.add_paragraph()
            p.text = item
            p.level = 0
            
            # Apply styling
            if font_size:
                p.font.size = Pt(font_size)
            if font_color and len(font_color) == 3:
                p.font.color.rgb = RGBColor(*font_color)
        
        # Right column
        right_frame = slide.placeholders[2].text_frame
        right_frame.clear()
        for item in right_content:
            p = right_frame.add_paragraph()
            p.text = item
            p.level = 0
            
            # Apply styling
            if font_size:
                p.font.size = Pt(font_size)
            if font_color and len(font_color) == 3:
                p.font.color.rgb = RGBColor(*font_color)
        
        # Apply borders if provided
        if border_color or border_width:
            for placeholder in [slide.placeholders[1], slide.placeholders[2]]:
                line = placeholder.line
                line.color.rgb = RGBColor(*border_color) if border_color and len(border_color) == 3 else RGBColor(0, 0, 0)
                line.width = Pt(border_width) if border_width else Pt(1.5)
        
        return {
            "success": True,
            "message": f"Added two-column slide: {title}",
            "slide_number": len(prs.slides)
        }
    
    def add_chart_slide(self, presentation_id: str, title: str, chart_type: str,
                       categories: List[str], series_data: List[Dict[str, Any]]) -> Dict[str, Any]:
        """
        Add a slide with a chart.
        
        Args:
            presentation_id: Presentation identifier
            title: Slide title
            chart_type: Type of chart (bar, column, line, pie)
            categories: Category labels
            series_data: List of series with name and values
            
        Returns:
            Dict with status and message
        """
        if presentation_id not in self.presentations:
            return {"success": False, "message": f"Presentation '{presentation_id}' not found"}
        
        prs = self.presentations[presentation_id]
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = title
        
        # Create chart data
        chart_data = CategoryChartData()
        chart_data.categories = categories
        
        for series in series_data:
            chart_data.add_series(series["name"], series["values"])
        
        # Map chart types
        chart_type_map = {
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE
        }
        
        # Add chart to slide
        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4.5)
        chart = slide.shapes.add_chart(
            chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED),
            x, y, cx, cy, chart_data
        ).chart
        
        return {
            "success": True,
            "message": f"Added {chart_type} chart slide: {title}",
            "slide_number": len(prs.slides),
            "chart_type": chart_type
        }
    
    def add_table_slide(self, presentation_id: str, title: str,
                       headers: List[str], rows: List[List[str]]) -> Dict[str, Any]:
        """
        Add a slide with a table.
        
        Args:
            presentation_id: Presentation identifier
            title: Slide title
            headers: Column headers
            rows: Table data rows
            
        Returns:
            Dict with status and message
        """
        if presentation_id not in self.presentations:
            return {"success": False, "message": f"Presentation '{presentation_id}' not found"}
        
        prs = self.presentations[presentation_id]
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = title
        
        # Calculate table dimensions
        num_cols = len(headers)
        num_rows = len(rows) + 1
        
        # Add table
        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4)
        table = slide.shapes.add_table(num_rows, num_cols, x, y, cx, cy).table
        
        # Set headers
        for col_idx, header in enumerate(headers):
            cell = table.rows[0].cells[col_idx]
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
        
        # Set data rows
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_value in enumerate(row_data):
                table.rows[row_idx + 1].cells[col_idx].text = str(cell_value)
        
        return {
            "success": True,
            "message": f"Added table slide: {title}",
            "slide_number": len(prs.slides),
            "table_size": f"{num_rows}x{num_cols}"
        }
    
    def add_image_slide(self, presentation_id: str, title: str, 
                       image_path: str, left: float = 1.0, top: float = 2.0,
                       width: Optional[float] = None, height: Optional[float] = None) -> Dict[str, Any]:
        """
        Add a slide with an image (diagrams, Gantt charts, flowcharts, etc.).
        
        Args:
            presentation_id: Presentation identifier
            title: Slide title
            image_path: Path to image file (PNG, JPG, SVG)
            left: Left position in inches (default: 1.0)
            top: Top position in inches (default: 2.0)
            width: Width in inches (optional, maintains aspect ratio if not provided)
            height: Height in inches (optional, maintains aspect ratio if not provided)
            
        Returns:
            Dict with status and message
        """
        if presentation_id not in self.presentations:
            return {"success": False, "message": f"Presentation '{presentation_id}' not found"}
        
        image_file = Path(image_path)
        if not image_file.exists():
            return {"success": False, "message": f"Image file not found: {image_path}"}
        
        prs = self.presentations[presentation_id]
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = title
        
        # Add image
        from pptx.util import Inches
        left_pos = Inches(left)
        top_pos = Inches(top)
        
        if width and height:
            width_size = Inches(width)
            height_size = Inches(height)
            slide.shapes.add_picture(str(image_file), left_pos, top_pos, width_size, height_size)
        elif width:
            width_size = Inches(width)
            slide.shapes.add_picture(str(image_file), left_pos, top_pos, width=width_size)
        elif height:
            height_size = Inches(height)
            slide.shapes.add_picture(str(image_file), left_pos, top_pos, height=height_size)
        else:
            # Auto-size to fit (max 8 inches wide)
            slide.shapes.add_picture(str(image_file), left_pos, top_pos, width=Inches(8))
        
        return {
            "success": True,
            "message": f"Added image slide: {title}",
            "slide_number": len(prs.slides),
            "image_path": str(image_file.absolute())
        }
    
    def add_gantt_chart_slide(self, presentation_id: str, title: str,
                             tasks: List[Dict[str, Any]]) -> Dict[str, Any]:
        """
        Add a slide with a Gantt chart (as a table representation).
        
        Args:
            presentation_id: Presentation identifier
            title: Slide title
            tasks: List of tasks with 'name', 'start', 'end', 'status' fields
            
        Returns:
            Dict with status and message
        """
        if presentation_id not in self.presentations:
            return {"success": False, "message": f"Presentation '{presentation_id}' not found"}
        
        prs = self.presentations[presentation_id]
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = title
        
        # Create Gantt-style table
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        
        num_cols = 5  # Task, Start, End, Duration, Status
        num_rows = len(tasks) + 1
        
        x, y, cx, cy = Inches(0.5), Inches(1.8), Inches(9), Inches(5)
        table = slide.shapes.add_table(num_rows, num_cols, x, y, cx, cy).table
        
        # Set column widths
        table.columns[0].width = Inches(3)  # Task name
        table.columns[1].width = Inches(1.5)  # Start
        table.columns[2].width = Inches(1.5)  # End
        table.columns[3].width = Inches(1.5)  # Duration
        table.columns[4].width = Inches(1.5)  # Status
        
        # Headers
        headers = ["Task", "Start Date", "End Date", "Duration", "Status"]
        for col_idx, header in enumerate(headers):
            cell = table.rows[0].cells[col_idx]
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            # Header background color
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Task rows
        for row_idx, task in enumerate(tasks):
            cells = table.rows[row_idx + 1].cells
            
            cells[0].text = task.get("name", "")
            cells[1].text = task.get("start", "")
            cells[2].text = task.get("end", "")
            cells[3].text = task.get("duration", "")
            cells[4].text = task.get("status", "")
            
            # Color code by status
            status = task.get("status", "").lower()
            if "complete" in status or "done" in status:
                cells[4].fill.solid()
                cells[4].fill.fore_color.rgb = RGBColor(146, 208, 80)
            elif "progress" in status or "active" in status:
                cells[4].fill.solid()
                cells[4].fill.fore_color.rgb = RGBColor(255, 217, 102)
            elif "blocked" in status or "delayed" in status:
                cells[4].fill.solid()
                cells[4].fill.fore_color.rgb = RGBColor(244, 176, 132)
            
            # Font size for all cells
            for cell in cells:
                cell.text_frame.paragraphs[0].font.size = Pt(10)
        
        return {
            "success": True,
            "message": f"Added Gantt chart slide: {title}",
            "slide_number": len(prs.slides),
            "task_count": len(tasks)
        }
    
    def add_process_flow_slide(self, presentation_id: str, title: str,
                               steps: List[str], flow_type: str = "horizontal") -> Dict[str, Any]:
        """
        Add a slide with process flow using chevron/arrow shapes.
        
        Args:
            presentation_id: Presentation identifier
            title: Slide title
            steps: List of process steps
            flow_type: "horizontal" or "vertical" flow direction
            
        Returns:
            Dict with status and message
        """
        if presentation_id not in self.presentations:
            return {"success": False, "message": f"Presentation '{presentation_id}' not found"}
        
        prs = self.presentations[presentation_id]
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = title
        
        # Create chevron/arrow shapes for process flow
        num_steps = len(steps)
        
        if flow_type == "horizontal":
            # Horizontal chevrons
            start_left = Inches(0.5)
            top = Inches(3)
            width = Inches(9.0 / num_steps)
            height = Inches(1.5)
            
            for idx, step in enumerate(steps):
                left = start_left + (width * idx)
                
                # Add chevron shape
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.CHEVRON,
                    left, top, width, height
                )
                
                # Set colors - alternate between two colors
                if idx % 2 == 0:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(68, 114, 196)
                else:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(112, 173, 71)
                
                # Add text
                text_frame = shape.text_frame
                text_frame.text = step
                text_frame.paragraphs[0].font.size = Pt(14)
                text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                text_frame.paragraphs[0].font.bold = True
                text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        else:  # vertical
            # Vertical arrows/boxes
            left = Inches(2)
            start_top = Inches(1.8)
            width = Inches(6)
            height = Inches(5.2 / num_steps)
            
            for idx, step in enumerate(steps):
                top = start_top + (height * idx)
                
                # Add rounded rectangle
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    left, top, width, height * 0.8
                )
                
                # Set colors
                if idx % 2 == 0:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(68, 114, 196)
                else:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(112, 173, 71)
                
                # Add text
                text_frame = shape.text_frame
                text_frame.text = f"{idx + 1}. {step}"
                text_frame.paragraphs[0].font.size = Pt(14)
                text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                text_frame.paragraphs[0].font.bold = True
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                
                # Add arrow between steps (except for last)
                if idx < num_steps - 1:
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.DOWN_ARROW,
                        left + width / 2 - Inches(0.25),
                        top + height * 0.8,
                        Inches(0.5),
                        height * 0.2
                    )
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = RGBColor(68, 114, 196)
                    arrow.line.color.rgb = RGBColor(68, 114, 196)
        
        return {
            "success": True,
            "message": f"Added process flow slide: {title}",
            "slide_number": len(prs.slides),
            "step_count": len(steps),
            "flow_type": flow_type
        }
    
    def add_timeline_slide(self, presentation_id: str, title: str,
                          events: List[Dict[str, str]]) -> Dict[str, Any]:
        """
        Add a slide with a timeline using shapes.
        
        Args:
            presentation_id: Presentation identifier
            title: Slide title
            events: List of events with 'date' and 'description' fields
            
        Returns:
            Dict with status and message
        """
        if presentation_id not in self.presentations:
            return {"success": False, "message": f"Presentation '{presentation_id}' not found"}
        
        prs = self.presentations[presentation_id]
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = title
        
        # Draw horizontal timeline
        timeline_top = Inches(3.5)
        timeline_left = Inches(1)
        timeline_width = Inches(8)
        
        # Timeline line
        line = slide.shapes.add_connector(
            1,  # Straight connector
            timeline_left, timeline_top,
            timeline_left + timeline_width, timeline_top
        )
        line.line.color.rgb = RGBColor(68, 114, 196)
        line.line.width = Pt(3)
        
        # Add events
        num_events = len(events)
        spacing = timeline_width / (num_events - 1) if num_events > 1 else Inches(0)
        
        for idx, event in enumerate(events):
            event_left = timeline_left + (spacing * idx)
            
            # Add circle marker
            marker = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                event_left - Inches(0.15),
                timeline_top - Inches(0.15),
                Inches(0.3),
                Inches(0.3)
            )
            marker.fill.solid()
            marker.fill.fore_color.rgb = RGBColor(68, 114, 196)
            marker.line.color.rgb = RGBColor(255, 255, 255)
            marker.line.width = Pt(2)
            
            # Add date above timeline
            date_box = slide.shapes.add_textbox(
                event_left - Inches(0.75),
                timeline_top - Inches(0.8),
                Inches(1.5),
                Inches(0.4)
            )
            date_frame = date_box.text_frame
            date_frame.text = event.get("date", "")
            date_frame.paragraphs[0].font.size = Pt(12)
            date_frame.paragraphs[0].font.bold = True
            date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Add description below timeline
            desc_box = slide.shapes.add_textbox(
                event_left - Inches(0.75),
                timeline_top + Inches(0.3),
                Inches(1.5),
                Inches(1)
            )
            desc_frame = desc_box.text_frame
            desc_frame.text = event.get("description", "")
            desc_frame.paragraphs[0].font.size = Pt(10)
            desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            desc_frame.word_wrap = True
        
        return {
            "success": True,
            "message": f"Added timeline slide: {title}",
            "slide_number": len(prs.slides),
            "event_count": len(events)
        }
    
    def add_diagram_slide(self, presentation_id: str, title: str,
                         diagram_type: str, items: List[str]) -> Dict[str, Any]:
        """
        Add a slide with a diagram (cycle, pyramid, matrix).
        
        Args:
            presentation_id: Presentation identifier
            title: Slide title
            diagram_type: "cycle", "pyramid", or "matrix"
            items: List of items to display
            
        Returns:
            Dict with status and message
        """
        if presentation_id not in self.presentations:
            return {"success": False, "message": f"Presentation '{presentation_id}' not found"}
        
        prs = self.presentations[presentation_id]
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = title
        
        if diagram_type == "cycle":
            # Create circular cycle diagram
            center_x = Inches(5)
            center_y = Inches(3.75)
            radius = Inches(2)
            num_items = len(items)
            
            import math
            for idx, item in enumerate(items):
                angle = (2 * math.pi * idx / num_items) - (math.pi / 2)
                x = center_x + radius * math.cos(angle) - Inches(0.6)
                y = center_y + radius * math.sin(angle) - Inches(0.4)
                
                # Add circle
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    x, y,
                    Inches(1.2),
                    Inches(0.8)
                )
                
                colors = [
                    RGBColor(68, 114, 196),
                    RGBColor(112, 173, 71),
                    RGBColor(255, 192, 0),
                    RGBColor(237, 125, 49)
                ]
                shape.fill.solid()
                shape.fill.fore_color.rgb = colors[idx % len(colors)]
                
                text_frame = shape.text_frame
                text_frame.text = item
                text_frame.paragraphs[0].font.size = Pt(11)
                text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                text_frame.paragraphs[0].font.bold = True
                text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                text_frame.word_wrap = True
        
        elif diagram_type == "pyramid":
            # Create pyramid diagram
            start_top = Inches(2)
            pyramid_height = Inches(5)
            max_width = Inches(7)
            center_x = Inches(5)
            
            level_height = pyramid_height / len(items)
            
            for idx, item in enumerate(items):
                # Width decreases as we go up
                width_ratio = 1 - (idx * 0.15)
                level_width = max_width * width_ratio
                
                left = center_x - (level_width / 2)
                top = start_top + (idx * level_height)
                
                # Add trapezoid
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.TRAPEZOID,
                    left, top,
                    level_width, level_height * 0.9
                )
                
                colors = [
                    RGBColor(68, 114, 196),
                    RGBColor(112, 173, 71),
                    RGBColor(255, 192, 0),
                    RGBColor(237, 125, 49)
                ]
                shape.fill.solid()
                shape.fill.fore_color.rgb = colors[idx % len(colors)]
                
                text_frame = shape.text_frame
                text_frame.text = item
                text_frame.paragraphs[0].font.size = Pt(14)
                text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                text_frame.paragraphs[0].font.bold = True
                text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        elif diagram_type == "matrix":
            # Create 2x2 matrix
            if len(items) != 4:
                return {"success": False, "message": "Matrix diagram requires exactly 4 items"}
            
            box_width = Inches(4)
            box_height = Inches(2.25)
            start_left = Inches(1)
            start_top = Inches(2)
            
            positions = [
                (start_left, start_top),  # Top-left
                (start_left + box_width, start_top),  # Top-right
                (start_left, start_top + box_height),  # Bottom-left
                (start_left + box_width, start_top + box_height)  # Bottom-right
            ]
            
            colors = [
                RGBColor(68, 114, 196),
                RGBColor(112, 173, 71),
                RGBColor(255, 192, 0),
                RGBColor(237, 125, 49)
            ]
            
            for idx, (left, top) in enumerate(positions):
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    left, top,
                    box_width * 0.95, box_height * 0.95
                )
                
                shape.fill.solid()
                shape.fill.fore_color.rgb = colors[idx]
                
                text_frame = shape.text_frame
                text_frame.text = items[idx]
                text_frame.paragraphs[0].font.size = Pt(14)
                text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                text_frame.paragraphs[0].font.bold = True
                text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                text_frame.word_wrap = True
        
        return {
            "success": True,
            "message": f"Added {diagram_type} diagram slide: {title}",
            "slide_number": len(prs.slides),
            "diagram_type": diagram_type
        }
    
    def add_shape_slide(self, presentation_id: str, title: str,
                       shapes: List[Dict[str, Any]]) -> Dict[str, Any]:
        """
        Add a slide with custom shapes (circles, squares, rectangles, etc.).
        
        Args:
            presentation_id: Presentation identifier
            title: Slide title
            shapes: List of shape definitions with:
                - shape_type: "circle", "square", "rectangle", "rounded_rectangle", 
                             "oval", "pentagon", "hexagon", "octagon", "triangle",
                             "diamond", "arrow", "star"
                - text: Text to display in shape (optional)
                - left: Left position in inches
                - top: Top position in inches
                - width: Width in inches (optional, auto-sized if text provided)
                - height: Height in inches (optional, auto-sized if text provided)
                - color: RGB color as [R, G, B] (optional, default: blue)
                - font_size: Font size in points (optional, default: 14)
                - text_color: RGB color for text as [R, G, B] (optional, default: white)
            
        Returns:
            Dict with status and message
        """
        if presentation_id not in self.presentations:
            return {"success": False, "message": f"Presentation '{presentation_id}' not found"}
        
        prs = self.presentations[presentation_id]
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = title
        
        # Map shape type names to MSO_SHAPE constants
        shape_type_map = {
            "circle": MSO_SHAPE.OVAL,
            "oval": MSO_SHAPE.OVAL,
            "square": MSO_SHAPE.RECTANGLE,
            "rectangle": MSO_SHAPE.RECTANGLE,
            "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
            "pentagon": MSO_SHAPE.REGULAR_PENTAGON,
            "hexagon": MSO_SHAPE.HEXAGON,
            "octagon": MSO_SHAPE.OCTAGON,
            "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
            "diamond": MSO_SHAPE.DIAMOND,
            "arrow": MSO_SHAPE.RIGHT_ARROW,
            "star": MSO_SHAPE.STAR_5_POINT,
            "star5": MSO_SHAPE.STAR_5_POINT,
            "star6": MSO_SHAPE.STAR_6_POINT,
            "star7": MSO_SHAPE.STAR_7_POINT,
            "cloud": MSO_SHAPE.CLOUD,
            "heart": MSO_SHAPE.HEART,
            "lightning": MSO_SHAPE.LIGHTNING_BOLT,
            "sun": MSO_SHAPE.SUN,
            "moon": MSO_SHAPE.MOON
        }
        
        added_shapes = []
        
        for shape_def in shapes:
            shape_type = shape_def.get("shape_type", "rectangle").lower()
            text = shape_def.get("text", "")
            left_pos = shape_def.get("left", 1.0)
            top_pos = shape_def.get("top", 2.0)
            
            # Get color (default: blue)
            color_rgb = shape_def.get("color", [68, 114, 196])
            fill_color = RGBColor(*color_rgb)
            
            # Get text color (default: white)
            text_color_rgb = shape_def.get("text_color", [255, 255, 255])
            text_color = RGBColor(*text_color_rgb)
            
            font_size = shape_def.get("font_size", 14)
            
            # Auto-size based on text if width/height not provided
            if text and ("width" not in shape_def or "height" not in shape_def):
                # Estimate size based on text length
                # Approximate: 0.1 inches per character width, 0.3 inches per line height
                text_length = len(text)
                estimated_width = max(1.5, min(6, text_length * 0.1))
                estimated_height = max(0.8, min(3, (text.count('\n') + 1) * 0.4))
                
                width = shape_def.get("width", estimated_width)
                height = shape_def.get("height", estimated_height)
                
                # For circles/ovals, make sure width and height are equal for perfect circles
                if shape_type == "circle":
                    dimension = max(width, height)
                    width = height = dimension
                elif shape_type == "square":
                    dimension = max(width, height)
                    width = height = dimension
            else:
                width = shape_def.get("width", 2.0)
                height = shape_def.get("height", 1.5)
                
                # Ensure square shapes have equal dimensions
                if shape_type == "square":
                    dimension = max(width, height)
                    width = height = dimension
                elif shape_type == "circle":
                    dimension = max(width, height)
                    width = height = dimension
            
            # Get MSO_SHAPE type
            mso_shape_type = shape_type_map.get(shape_type, MSO_SHAPE.RECTANGLE)
            
            # Add shape to slide
            shape = slide.shapes.add_shape(
                mso_shape_type,
                Inches(left_pos),
                Inches(top_pos),
                Inches(width),
                Inches(height)
            )
            
            # Set fill color
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
            
            # Add text if provided
            if text:
                text_frame = shape.text_frame
                text_frame.clear()
                text_frame.text = text
                
                # Format text
                paragraph = text_frame.paragraphs[0]
                paragraph.font.size = Pt(font_size)
                paragraph.font.color.rgb = text_color
                paragraph.font.bold = True
                paragraph.alignment = PP_ALIGN.CENTER
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                text_frame.word_wrap = True
            
            added_shapes.append({
                "type": shape_type,
                "text": text,
                "position": f"({left_pos}, {top_pos})",
                "size": f"{width}x{height}"
            })
        
        return {
            "success": True,
            "message": f"Added shape slide with {len(shapes)} shape(s): {title}",
            "slide_number": len(prs.slides),
            "shapes": added_shapes
        }
    
    def save_presentation(self, presentation_id: str, output_path: str) -> Dict[str, Any]:
        """
        Save the presentation to a file with automatic timestamp in filename.
        
        Args:
            presentation_id: Presentation identifier
            output_path: Output file path (timestamp will be added before .pptx extension)
            
        Returns:
            Dict with status and message including the actual file path with timestamp
        """
        if presentation_id not in self.presentations:
            return {"success": False, "message": f"Presentation '{presentation_id}' not found"}
        
        prs = self.presentations[presentation_id]
        
        # Create directory if it doesn't exist
        output_file = Path(output_path)
        output_file.parent.mkdir(parents=True, exist_ok=True)
        
        # Add timestamp to filename (same format as charts: YYYYMMDD_HHMMSS)
        ts = datetime.now().strftime('%Y%m%d_%H%M%S')
        stem = output_file.stem  # filename without extension
        suffix = output_file.suffix  # extension (e.g., .pptx)
        timestamped_name = f"{stem}_{ts}{suffix}"
        timestamped_path = output_file.parent / timestamped_name
        
        # Save presentation with timestamped filename
        prs.save(str(timestamped_path))
        
        return {
            "success": True,
            "message": f"Saved presentation to {timestamped_path}",
            "slide_count": len(prs.slides),
            "file_path": str(timestamped_path.absolute())
        }
    
    def list_presentations(self) -> Dict[str, Any]:
        """
        List all active presentations.
        
        Returns:
            Dict with presentation information
        """
        if not self.presentations:
            return {
                "success": True,
                "message": "No active presentations",
                "presentations": []
            }
        
        presentation_list = []
        for pres_id, prs in self.presentations.items():
            presentation_list.append({
                "id": pres_id,
                "slide_count": len(prs.slides)
            })
        
        return {
            "success": True,
            "message": f"Found {len(self.presentations)} active presentation(s)",
            "presentations": presentation_list
        }
    
    def get_presentation(self, presentation_id: str) -> Optional[Presentation]:
        """Get a presentation by ID."""
        return self.presentations.get(presentation_id)
    
    def delete_presentation(self, presentation_id: str) -> Dict[str, Any]:
        """Delete a presentation from memory."""
        if presentation_id not in self.presentations:
            return {"success": False, "message": f"Presentation '{presentation_id}' not found"}
        
        del self.presentations[presentation_id]
        return {
            "success": True,
            "message": f"Deleted presentation '{presentation_id}'"
        }
